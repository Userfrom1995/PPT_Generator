from transformers import AutoTokenizer, AutoModelForCausalLM
from pptx import Presentation
import torch
import gc

# Load GPT-J model and move it to the GPU
model_name = "EleutherAI/gpt-j-6B"
tokenizer = AutoTokenizer.from_pretrained(model_name)
model = AutoModelForCausalLM.from_pretrained(model_name)

# Set padding token to be the same as the EOS token if not already set
if tokenizer.pad_token is None:
    tokenizer.pad_token = tokenizer.eos_token

# Generate text from description using GPU and attention_mask
def generate_text(description):
    # Tokenize the input and generate attention mask
    inputs = tokenizer(description, return_tensors="pt", padding=True, truncation=True)
    input_ids = inputs.input_ids
    attention_mask = inputs.attention_mask
    
    # Generate text using the model on the GPU
    with torch.no_grad():  # Disable gradient calculation for inference
        output = model.generate(input_ids, attention_mask=attention_mask, max_length=150, num_return_sequences=1)
    
    # Clear GPU memory
    del input_ids
    del attention_mask
    del output
    torch.cuda.empty_cache()
    
    # Run garbage collection
    gc.collect()
    
    return tokenizer.decode(output[0], skip_special_tokens=True)

# Create PowerPoint from generated text
def create_ppt_from_text(text, output_file="generated_presentation.pptx"):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "AI Generated Slide"
    content = slide.placeholders[1]
    content.text = text
    prs.save(output_file)
    print(f"Presentation saved as {output_file}")

# Example Usage
description = "Create a presentation about the latest advancements in AI technology."
generated_text = generate_text(description)
create_ppt_from_text(generated_text)
